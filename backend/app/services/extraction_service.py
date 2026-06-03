import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from typing import Dict, Any, Tuple
import re

class ExtractionService:
    @staticmethod
    def detect_topic(text: str) -> str:
        topics_keywords = {
            "Data Structures & Algorithms": ["dsa", "array", "list", "stack", "queue", "tree", "graph", "sorting", "searching", "complexity", "big o", "recursion"],
            "Full-Stack Development": ["html", "css", "javascript", "react", "node", "express", "backend", "frontend", "api", "database", "mongodb", "fastapi", "rest"],
            "Cybersecurity & Cryptography": ["security", "cyber", "cryptography", "encryption", "decryption", "hashing", "network", "firewall", "vulnerability", "malware", "owasp", "penetration"],
            "Artificial Intelligence & ML": ["ai", "machine learning", "neural network", "deep learning", "supervised", "unsupervised", "regression", "classification", "nlp", "computer vision", "dataset"],
            "Database Management": ["sql", "nosql", "query", "index", "normalization", "transaction", "acid", "schema", "postgres", "mysql", "mongodb"]
        }
        
        scores = {topic: 0 for topic in topics_keywords}
        text_lower = text.lower()
        
        for topic, keywords in topics_keywords.items():
            for word in keywords:
                # Count occurrences of keywords
                occurrences = len(re.findall(r'\b' + re.escape(word) + r'\b', text_lower))
                scores[topic] += occurrences
                
        best_topic = max(scores, key=scores.get)
        if scores[best_topic] > 0:
            return best_topic
        return "General Study Material"

    @staticmethod
    def extract_pdf(file_path: str) -> Tuple[str, int]:
        text = ""
        doc = fitz.open(file_path)
        page_count = len(doc)
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        return text, page_count

    @staticmethod
    def extract_docx(file_path: str) -> Tuple[str, int]:
        text = ""
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        word_count = len(text.split())
        page_count = max(1, round(word_count / 350))
        return text, page_count

    @staticmethod
    def extract_pptx(file_path: str) -> Tuple[str, int]:
        text = ""
        prs = Presentation(file_path)
        page_count = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            text += f"[Slide {i+1}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text, page_count

    @staticmethod
    def extract_txt(file_path: str) -> Tuple[str, int]:
        text = ""
        for encoding in ["utf-8", "latin-1"]:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                break
            except Exception:
                continue
        word_count = len(text.split())
        page_count = max(1, round(word_count / 400))
        return text, page_count

    @classmethod
    def extract_metadata_and_text(cls, file_path: str, extension: str) -> Dict[str, Any]:
        ext = extension.lower().strip(".")
        
        if ext == "pdf":
            text, page_count = cls.extract_pdf(file_path)
        elif ext == "docx":
            text, page_count = cls.extract_docx(file_path)
        elif ext == "pptx":
            text, page_count = cls.extract_pptx(file_path)
        elif ext == "txt":
            text, page_count = cls.extract_txt(file_path)
        else:
            raise ValueError(f"Unsupported file extension: .{ext}")

        word_count = len(text.split())
        estimated_read_time = max(1, round(word_count / 200))
        detected_topic = cls.detect_topic(text)
        
        return {
            "text": text,
            "word_count": word_count,
            "page_count": page_count,
            "estimated_read_time": estimated_read_time,
            "detected_topic": detected_topic
        }
